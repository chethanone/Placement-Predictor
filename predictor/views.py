from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from django.http import JsonResponse
from django.core.files.storage import FileSystemStorage
from django.utils import timezone
import json
import random
import re
import os
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()

# Import Google Gemini AI
try:
    import google.generativeai as genai
    # Configure Gemini API from environment variable
    GEMINI_API_KEY = os.getenv('GEMINI_API_KEY')
    if GEMINI_API_KEY:
        genai.configure(api_key=GEMINI_API_KEY)
    else:
        print(" WARNING: GEMINI_API_KEY not found in .env file")
        genai = None
except ImportError:
    genai = None
    print(" WARNING: google-generativeai not installed")

# Import PDF/PPT processing libraries
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

from .models import (
    StudentRecord, 
    StudentQuiz, 
    QuizQuestion, 
    StudentPrediction,
    TrainingSession,
    SessionRecommendation
)

def process_pdf(file_path):
    text = ""
    print(f" Processing PDF: {file_path}")
    
    # Try PyMuPDF first (best quality)
    if fitz:
        try:
            print(" Using PyMuPDF for high-quality extraction...")
            doc = fitz.open(file_path)
            for page_num, page in enumerate(doc, 1):
                # Extract text with better formatting
                page_text = page.get_text("text", sort=True)  # Sort text by reading order
                if page_text.strip():
                    # Don't add page markers - just append text with spacing
                    text += "\n" + page_text
                else:
                    # Try blocks method if normal extraction fails
                    blocks = page.get_text("blocks")
                    for block in blocks:
                        if len(block) >= 5:  # Block has text
                            text += block[4] + " "
            doc.close()
            
            # Clean extracted text - Remove junk characters and excessive whitespace
            import re
            
            # Remove ALL page markers in various formats
            text = re.sub(r'---\s*[Pp]age\s+\d+\s*---', '', text)  # --- Page 17 ---
            text = re.sub(r'-+\s*[Pp]age\s+\d+\s*-+', '', text)   # -- page 17 --
            text = re.sub(r'[Pp]age\s+\d+\s+of\s+\d+', '', text)   # Page 1 of 10
            text = re.sub(r'^\s*\d+\s*$', '', text, flags=re.MULTILINE)  # Remove standalone numbers
            
            # Remove special characters and control characters (but keep newlines, spaces, basic punctuation)
            text = re.sub(r'[\x00-\x08\x0b-\x0c\x0e-\x1f\x7f-\x9f]', '', text)  # Remove control chars
            
            # Remove excessive blank lines (more than 2 consecutive)
            text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)
            
            # Remove excessive spaces
            text = re.sub(r' {2,}', ' ', text)
            
            # Remove bullet points and list markers if they're standalone
            text = re.sub(r'^\s*[•○▪▫■□●◆◇★☆]+\s*$', '', text, flags=re.MULTILINE)
            
            # Clean up common PDF artifacts
            text = re.sub(r'[\uf0b7\uf0a7\uf076\uf0d8]', '', text)  # Remove common PDF bullet symbols
            
            # Final cleanup - trim leading/trailing whitespace
            text = text.strip()
            
            print(f" Extracted and cleaned {len(text)} characters with PyMuPDF")
            if len(text) > 100:
                return text
        except Exception as e:
            print(f" PyMuPDF error: {e}")
    
    # Fallback to PyPDF2
    if PyPDF2 and len(text) < 100:
        try:
            print(" Using PyPDF2 fallback...")
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(reader.pages, 1):
                    page_text = page.extract_text()
                    if page_text.strip():
                        # Don't add page markers - just append text
                        text += "\n" + page_text
            
            # Clean text with same cleaning as above
            import re
            text = re.sub(r'---\s*[Pp]age\s+\d+\s*---', '', text)
            text = re.sub(r'-+\s*[Pp]age\s+\d+\s*-+', '', text)
            text = re.sub(r'[Pp]age\s+\d+\s+of\s+\d+', '', text)
            text = re.sub(r'[\x00-\x08\x0b-\x0c\x0e-\x1f\x7f-\x9f]', '', text)
            text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)
            text = re.sub(r' {2,}', ' ', text)
            text = re.sub(r'[\uf0b7\uf0a7\uf076\uf0d8]', '', text)
            text = text.strip()
            
            print(f" Extracted and cleaned {len(text)} characters with PyPDF2")
        except Exception as e:
            print(f" PyPDF2 error: {e}")
    
    if len(text) < 100:
        print(" WARNING: Very little text extracted - PDF may be image-based or encrypted")
    
    return text

def process_ppt(file_path):
    text = ""
    print(f" Processing PPT: {file_path}")
    if Presentation:
        try:
            print(" Using python-pptx...")
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            print(f" Extracted {len(text)} characters")
        except Exception as e:
            print(f" PPT error: {e}")
            pass
    else:
        print(" python-pptx not available")
    return text

def process_uploaded_file(file_path):
    print(f"\n Processing uploaded file: {file_path}")
    ext = os.path.splitext(file_path)[1].lower()
    print(f" Extension: {ext}")
    if ext == '.pdf':
        return process_pdf(file_path)
    elif ext in ['.ppt', '.pptx']:
        return process_ppt(file_path)
    return ""

def generate_questions_from_text(text, max_questions=25):
    """Generate 25 diverse quiz questions using Google Gemini AI - OPTIMIZED"""
    if not text:
        print(" No text provided for question generation")
        return []
    
    # Check if Gemini is available
    if not genai:
        print(" Gemini AI not available, using fallback")
        return generate_basic_questions(text, max_questions)
    
    try:
        print(f" Using Gemini AI (FAST MODE) to generate {max_questions} questions...")
        print(f" Content length: {len(text)} characters")
        
        # Initialize Gemini model with faster generation config
        model = genai.GenerativeModel(
            'gemini-2.0-flash',
            generation_config={
                'temperature': 0.7,
                'top_p': 0.95,
                'top_k': 40,
                'max_output_tokens': 4096,
            }
        )
        
        # Create optimized prompt for faster generation with better question quality
        prompt = f"""
Analyze this educational content and generate EXACTLY 25 high-quality quiz questions that test understanding of KEY CONCEPTS and IMPORTANT TOPICS.

CONTENT TO ANALYZE:
{text[:15000]}

INSTRUCTIONS FOR QUESTION GENERATION:
1. Focus on MAIN CONCEPTS, definitions, algorithms, formulas, and key principles
2. Avoid trivial questions about formatting, page numbers, or irrelevant details
3. Make MCQ options plausible and challenging (not obviously wrong)
4. For fill-in-blank, choose important terms/concepts (not random words)
5. For true/false, test understanding of core principles
6. For short answers, ask about explanations or applications of concepts
7. Extract a relevant text snippet for the "ref" field that explains the answer

OUTPUT FORMAT (STRICT JSON, NO MARKDOWN):
{{
  "questions": [
    {{"q": "What is the primary purpose of divide and conquer?", "type": "mcq", "options": ["Break problems into subproblems", "Sort arrays", "Search trees", "Optimize space"], "answer": "Break problems into subproblems", "ref": "Divide and conquer works by recursively breaking down a problem into sub-problems until they become simple enough to solve directly."}},
    {{"q": "The time complexity of binary search is _____.", "type": "fill_blank", "answer": "O(log n)", "ref": "Binary search has logarithmic time complexity O(log n) because it divides the search space in half with each comparison."}},
    {{"q": "Dynamic programming stores results of subproblems to avoid recomputation.", "type": "true_false", "answer": "True", "ref": "Dynamic programming is an optimization technique that stores the results of expensive function calls and reuses them when the same inputs occur again."}},
    {{"q": "Explain the greedy approach in algorithm design.", "type": "short_answer", "answer": "Makes locally optimal choices at each step hoping to find global optimum", "ref": "A greedy algorithm makes the locally optimal choice at each stage with the hope of finding a global optimum."}}
  ]
}}

DISTRIBUTION REQUIRED:
- 10 Multiple Choice Questions (MCQ) with 4 meaningful options each
- 5 Fill in the Blank questions about key terms/concepts
- 5 True/False statements about core principles  
- 5 Short Answer questions requiring explanation
- Mix the order randomly (don't group by type)
"""
        
        # Generate content with Gemini (faster with optimized prompt)
        print("⏳ Calling Gemini API...")
        response = model.generate_content(prompt)
        response_text = response.text.strip()
        print(f" Received response ({len(response_text)} chars)")
        
        # Clean the response
        if response_text.startswith('```json'):
            response_text = response_text[7:]
        if response_text.startswith('```'):
            response_text = response_text[3:]
        if response_text.endswith('```'):
            response_text = response_text[:-3]
        response_text = response_text.strip()
        
        print(" Parsing JSON response...")
        # Parse JSON response
        result = json.loads(response_text)
        questions = result.get('questions', [])
        print(f" Parsed {len(questions)} questions from JSON")
        
        # Randomize the order of questions (shuffle them)
        random.shuffle(questions)
        print(" Questions randomized")
        
        # Convert to our format (optimized mapping)
        formatted_questions = []
        for q in questions[:max_questions]:
            formatted_questions.append({
                'question_text': q.get('q', q.get('question_text', '')),
                'correct_answer': q.get('answer', q.get('correct_answer', '')),
                'question_type': q.get('type', q.get('question_type', 'text')),
                'options': q.get('options', []),
                'reference_text': q.get('ref', q.get('reference_text', ''))
            })
        
        # Ensure we have exactly 25 questions
        while len(formatted_questions) < max_questions:
            formatted_questions.append({
                'question_text': f"Question {len(formatted_questions) + 1} based on content.",
                'correct_answer': "Answer based on content",
                'question_type': 'short_answer',
                'options': []
            })
        
        print(f" Returning {len(formatted_questions)} formatted questions (randomized)")
        return formatted_questions[:max_questions]
        
    except json.JSONDecodeError as e:
        print(f" JSON Parse Error: {str(e)}")
        return generate_basic_questions(text, max_questions)
    except Exception as e:
        print(f" Gemini API Error: {str(e)}")
        import traceback
        traceback.print_exc()
        return generate_basic_questions(text, max_questions)

def generate_basic_questions(text, max_questions=25):
    """Fallback: Generate basic questions if Gemini API fails."""
    if not text:
        return []
    
    # Clean text and split into sentences
    sentences = re.split(r'(?<=[.!?])\s+', text)
    
    # Filter out very short sentences, page markers, and figure references
    good_sentences = []
    for s in sentences:
        s = s.strip()
        # Skip if too short, is a page marker, or is just a figure reference
        if (len(s) < 30 or 
            re.match(r'^(Page|\d+|Figure|Fig\.|Table|Diagram)', s, re.IGNORECASE) or
            s.count(' ') < 3):  # At least 3 words
            continue
        good_sentences.append(s)
    
    if not good_sentences:
        # If no good sentences found, return minimal questions
        return [{
            'question_text': "Based on the content provided, what are the key concepts discussed?",
            'correct_answer': "Please refer to the study material",
            'question_type': 'short_answer',
            'options': [],
            'reference_text': "Please review the uploaded document for details."
        }] * max_questions
    
    questions = []
    random.shuffle(good_sentences)
    
    question_templates = [
        ('mcq', lambda s: (f"Which statement best describes the following?\n{s[:200]}...", 
                          "Correct", 
                          ["Correct", "Incorrect", "Partially correct", "Cannot determine"])),
        ('fill_blank', lambda s: create_fill_blank(s)),
        ('true_false', lambda s: (f"The following statement is true: {s[:150]}...", 
                                  "True", 
                                  ["True", "False"])),
        ('short_answer', lambda s: (f"Explain the concept: {s[:120]}...", 
                                   s[:200], 
                                   []))
    ]
    
    for i, sentence in enumerate(good_sentences[:max_questions * 2]):  # Get extra in case some fail
        if len(questions) >= max_questions:
            break
            
        template_type, template_func = question_templates[i % 4]
        
        try:
            # Check if sentence mentions figures/diagrams
            has_figure = bool(re.search(r'\b(figure|fig\.|diagram|graph|chart|table|image)\b', sentence, re.IGNORECASE))
            reference_note = f"[Note: This content may reference visual elements from the original document]\n\n{sentence}" if has_figure else sentence
            
            if template_type == 'fill_blank':
                result = template_func(sentence)
                if result:  # Only add if fill_blank was successful
                    q_text, answer = result
                    questions.append({
                        'question_text': q_text,
                        'correct_answer': answer,
                        'question_type': 'fill_blank',
                        'options': [],
                        'reference_text': reference_note
                    })
            else:
                q_text, answer, options = template_func(sentence)
                questions.append({
                    'question_text': q_text,
                    'correct_answer': answer,
                    'question_type': template_type,
                    'options': options,
                    'reference_text': reference_note
                })
        except Exception as e:
            print(f" Error creating question: {e}")
            continue
    
    # If still not enough questions, add generic ones
    while len(questions) < max_questions:
        questions.append({
            'question_text': f"What are the key takeaways from the study material? (Topic {len(questions) + 1})",
            'correct_answer': "Refer to the study material for comprehensive understanding",
            'question_type': 'short_answer',
            'options': [],
            'reference_text': "This question requires comprehensive review of the uploaded document."
        })
    
    return questions[:max_questions]

def create_fill_blank(sentence):
    """Helper function to create fill-in-the-blank questions."""
    words = sentence.split()
    if len(words) < 6:
        return None
    
    # Find important words (nouns, longer words, capitalized words)
    important_words = []
    for i, word in enumerate(words):
        # Clean word of punctuation
        clean_word = re.sub(r'[^\w]', '', word)
        if (len(clean_word) > 5 and 
            clean_word.isalnum() and 
            not clean_word.isdigit() and
            i > 0 and i < len(words) - 1):  # Not first or last word
            important_words.append((i, word, clean_word))
    
    if not important_words:
        return None
    
    # Pick a random important word
    idx, original_word, clean_word = random.choice(important_words)
    question_text = ' '.join(words[:idx] + ['_____'] + words[idx+1:])
    
    # Limit question length
    if len(question_text) > 200:
        question_text = question_text[:200] + '...'
    
    return (question_text, clean_word)

def student_quiz_upload(request):
    if request.method == 'POST':
        try:
            # Get student_id from session first, fallback to POST data
            student_id = request.session.get('student_usn') or request.POST.get('student_id')
            student_name = request.session.get('student_name') or request.POST.get('student_name') or student_id
            uploaded_file = request.FILES.get('file')
            
            print(f" Quiz Upload Request - Student: {student_id}, Name: {student_name}")
            
            if not student_id:
                messages.error(request, "Session expired. Please log in again.")
                return redirect('student_entry')
            
            if not uploaded_file:
                messages.error(request, "Please upload a file.")
                return redirect('student_quiz_upload')
            
            print(f" File uploaded: {uploaded_file.name}")
            
            # Save file
            fs = FileSystemStorage(location='media/syllabi')
            filename = fs.save(f"{student_id}_{uploaded_file.name}", uploaded_file)
            file_path = fs.path(filename)
            
            print(f" File saved to: {file_path}")
            
            # Process file
            print(" Extracting text from file...")
            syllabus_text = process_uploaded_file(file_path)
            
            if not syllabus_text:
                messages.error(request, "Could not extract text from file. Please ensure it's a valid PDF or PPT.")
                print(" Text extraction failed")
                return redirect('student_quiz_upload')
            
            print(f" Text extracted: {len(syllabus_text)} characters")
            
            # Create quiz
            print(" Creating quiz record...")
            
            # Determine file type
            file_ext = os.path.splitext(uploaded_file.name)[1].lower().replace('.', '')
            
            quiz = StudentQuiz.objects.create(
                student_id=student_id,
                student_name=student_name,
                file_type=file_ext,
                extracted_text=syllabus_text[:50000]  # Store more text for better question generation (50KB)
            )
            
            # Save the uploaded file to the quiz
            quiz.uploaded_file.save(filename, uploaded_file, save=True)
            
            print(f" Quiz created with ID: {quiz.id}")
            
            # Generate and save questions
            print(" Generating questions with AI...")
            questions = generate_questions_from_text(syllabus_text)
            print(f" Generated {len(questions)} questions")
            
            # Save questions with proper numbering
            question_type_count = {}
            for idx, q in enumerate(questions, 1):
                q_type = q.get('question_type', 'text')
                question_type_count[q_type] = question_type_count.get(q_type, 0) + 1
                
                QuizQuestion.objects.create(
                    quiz=quiz,
                    question_number=idx,
                    question_text=q['question_text'],
                    correct_answer=q['correct_answer'],
                    question_type=q_type,
                    options=json.dumps(q.get('options', [])) if q.get('options') else None,
                    reference_text=q.get('reference_text', '')
                )
                print(f" Question {idx} ({q_type}) saved")
            
            print(f"\n Question Type Distribution:")
            for q_type, count in question_type_count.items():
                print(f" {q_type}: {count} questions")
            
            quiz.questions_generated = True
            quiz.save()
            
            print(f" Quiz generation complete! Redirecting to quiz {quiz.id}")
            # Removed success message - user will see the quiz directly
            return redirect('student_take_quiz', quiz_id=quiz.id)
            
        except Exception as e:
            print(f" ERROR in quiz upload: {str(e)}")
            import traceback
            traceback.print_exc()
            messages.error(request, f"Error generating quiz: {str(e)}")
            return redirect('student_quiz_upload')
    
    return render(request, 'predictor/student/quiz_upload.html')

def student_take_quiz(request, quiz_id):
    # Clear any previous error messages from upload
    storage = messages.get_messages(request)
    storage.used = True
    
    quiz = get_object_or_404(StudentQuiz, id=quiz_id)
    # Get questions ordered by question_number (which preserves randomization from generation)
    questions = QuizQuestion.objects.filter(quiz=quiz).order_by('question_number')
    return render(request, 'predictor/student/take_quiz.html', {
        'quiz': quiz,
        'questions': questions
    })

def student_retake_quiz(request, quiz_id):
    """
    Retake the same quiz with randomized questions without re-uploading the file.
    This regenerates questions from the already extracted text.
    """
    quiz = get_object_or_404(StudentQuiz, id=quiz_id)
    
    try:
        # Delete old questions
        QuizQuestion.objects.filter(quiz=quiz).delete()
        
        # Regenerate questions from stored text
        print(f" Regenerating questions for quiz {quiz_id}...")
        syllabus_text = quiz.extracted_text
        
        if not syllabus_text:
            messages.error(request, "No text available for this quiz. Please upload the file again.")
            return redirect('student_quiz_upload')
        
        # Generate new randomized questions
        questions = generate_questions_from_text(syllabus_text)
        print(f" Generated {len(questions)} new randomized questions")
        
        # Save new questions
        question_type_count = {}
        for idx, q in enumerate(questions, 1):
            q_type = q.get('question_type', 'text')
            question_type_count[q_type] = question_type_count.get(q_type, 0) + 1
            
            QuizQuestion.objects.create(
                quiz=quiz,
                question_number=idx,
                question_text=q['question_text'],
                correct_answer=q['correct_answer'],
                question_type=q_type,
                options=json.dumps(q.get('options', [])) if q.get('options') else None,
                reference_text=q.get('reference_text', '')
            )
        
        # Reset quiz metadata
        quiz.submitted = False
        quiz.score = None
        quiz.start_time = timezone.now()
        quiz.end_time = None
        quiz.save()
        
        print(f" Quiz {quiz_id} reset successfully!")
        messages.success(request, "🔄 Quiz regenerated with new randomized questions!")
        return redirect('student_take_quiz', quiz_id=quiz.id)
        
    except Exception as e:
        print(f" ERROR retaking quiz: {str(e)}")
        messages.error(request, f"Error regenerating quiz: {str(e)}")
        return redirect('student_quiz_upload')

def verify_answer_with_content(question, student_answer, pdf_content):
    """
    FAST & ACCURATE answer verification using AI and content matching.
    Optimized for speed while maintaining accuracy.
    """
    if not student_answer or not student_answer.strip():
        return False, "No answer provided"
    
    student_answer = student_answer.strip()
    correct_answer = question.correct_answer.strip()
    question_text = question.question_text
    
    # For MCQ and True/False, use exact matching (FASTEST)
    if question.question_type == 'mcq':
        is_correct = (student_answer == correct_answer)
        return is_correct, "Exact match required for MCQ"
    
    if question.question_type in ['true_false', 'tf']:
        is_correct = (student_answer.lower() == correct_answer.lower())
        return is_correct, "Exact match required for True/False"
    
    # For Fill in the Blank and Short Answer, use optimized AI verification
    if question.question_type in ['fill_blank', 'short_answer']:
        # Quick exact match check first (FAST PATH)
        correct_lower = correct_answer.lower()
        answer_lower = student_answer.lower()
        
        # Remove punctuation and extra spaces for flexible matching
        import string
        translator = str.maketrans('', '', string.punctuation)
        correct_normalized = correct_lower.translate(translator).strip()
        answer_normalized = answer_lower.translate(translator).strip()
        
        # Exact match after normalization (best case)
        if correct_normalized == answer_normalized:
            return True, "Answer matches expected response"
        
        # Basic keyword matching for fill_blank with STRICT comparison
        if question.question_type == 'fill_blank':
            # Split into words and filter out very short words
            correct_words = [w for w in correct_normalized.split() if len(w) > 2]
            answer_words = [w for w in answer_normalized.split() if len(w) > 2]
            
            if not correct_words:
                # If expected answer is very short, require exact match
                return correct_normalized == answer_normalized, "Exact match required"
            
            # Check if answer contains the expected answer as substring (flexible)
            if correct_normalized in answer_normalized or answer_normalized in correct_normalized:
                return True, "Answer contains expected response"
            
            # Check word overlap - require HIGH similarity (at least 80% of words must match)
            correct_words_set = set(correct_words)
            answer_words_set = set(answer_words)
            common_words = correct_words_set & answer_words_set
            
            # Calculate match ratio
            if len(correct_words_set) > 0:
                match_ratio = len(common_words) / len(correct_words_set)
                # Require at least 80% word match AND at least one word matches
                if match_ratio >= 0.8 and len(common_words) > 0:
                    return True, f"High similarity ({match_ratio:.0%} words match)"
            
            # If no good match, return False BEFORE AI check for fill_blank
            # This prevents "random text" from being marked correct
            if len(common_words) == 0:
                return False, "Answer does not match expected response"
        
        # Fast AI verification if Gemini available (OPTIMIZED)
        if genai and GEMINI_API_KEY and pdf_content:
            try:
                model = genai.GenerativeModel(
                    'gemini-2.0-flash',
                    generation_config={
                        'temperature': 0.2,  # Lower for more consistent answers
                        'max_output_tokens': 150,  # Reduced for faster response
                    }
                )
                
                # Concise verification prompt (OPTIMIZED FOR SPEED)
                verification_prompt = f"""
Content: {pdf_content[:2000]}

Q: {question_text}
Expected Answer: {correct_answer}
Student Answer: {student_answer}

Verify if student answer is SEMANTICALLY CORRECT. Ignore punctuation and case, but the MEANING must match. Be STRICT - if the answer is wrong or random text, mark it incorrect.

JSON only:
{{"correct": true/false, "reason": "brief"}}
"""
                
                response = model.generate_content(verification_prompt)
                response_text = response.text.strip().replace('```json', '').replace('```', '').strip()
                
                result = json.loads(response_text)
                return result.get('correct', result.get('is_correct', False)), result.get('reason', result.get('reasoning', 'AI verification'))
                
            except Exception as e:
                print(f" AI verification failed: {e}")
                # Fall through to keyword matching
        
        # Fallback keyword matching for short answer (FAST)
        if question.question_type == 'short_answer':
            correct_keywords = set(word.lower() for word in correct_answer.split() if len(word) > 3)
            answer_keywords = set(word.lower() for word in student_answer.split() if len(word) > 3)
            
            if not correct_keywords:
                return len(student_answer) > 10, "Length check"
            
            common_keywords = correct_keywords & answer_keywords
            match_ratio = len(common_keywords) / len(correct_keywords)
            
            if match_ratio >= 0.4:  # Lowered threshold for better accuracy
                return True, f"Match: {match_ratio*100:.0f}%"
            else:
                return False, f"Low match: {match_ratio*100:.0f}%"
    
    # Default fallback
    return correct_answer.lower() in student_answer.lower(), "Default matching"


def student_submit_quiz(request, quiz_id):
    """Enhanced quiz submission with AI-powered answer verification"""
    if request.method == 'POST':
        quiz = get_object_or_404(StudentQuiz, id=quiz_id)
        questions = QuizQuestion.objects.filter(quiz=quiz)
        score = 0
        total = questions.count()
        
        # Get time taken
        time_taken = int(request.POST.get('time_taken', 0))
        
        print(f" Submitting quiz {quiz_id} with {total} questions")
        print(f" PDF Content available: {len(quiz.extracted_text)} characters")
        
        for question in questions:
            answer_key = f'question_{question.id}'
            answer = request.POST.get(answer_key, '').strip()
            
            print(f" Q{question.question_number}: Type={question.question_type}, Answer='{answer[:50]}'")
            
            if answer:
                question.student_answer = answer
                
                # Use enhanced verification with PDF content
                is_correct, reasoning = verify_answer_with_content(
                    question, 
                    answer, 
                    quiz.extracted_text
                )
                
                question.is_correct = is_correct
                question.save()
                
                if is_correct:
                    score += 1
                    print(f" Correct! ({reasoning})")
                else:
                    print(f" Incorrect ({reasoning})")
                    print(f" Expected: {question.correct_answer[:50]}")
            else:
                question.is_correct = False
                question.save()
                print(f" No answer provided")
        
        # Calculate results
        percentage = (score / total * 100) if total > 0 else 0
        quiz.score = score
        quiz.percentage = percentage
        quiz.status = 'completed'
        
        # Store time taken in minutes
        quiz.completed_at = timezone.now()
        quiz.save()
        
        print(f" Quiz submitted: {score}/{total} ({percentage:.1f}%)")
        # Removed success message - user will see results directly
        
        return redirect('student_quiz_results', quiz_id=quiz.id)
    
    return redirect('student_take_quiz', quiz_id=quiz_id)


def get_study_recommendations(weak_areas, percentage, quiz_content):
    """Generate FAST AI-powered study recommendations based on uploaded PDF content"""
    recommendations = []
    
    if not genai or not GEMINI_API_KEY:
        return recommendations
    
    try:
        model = genai.GenerativeModel(
            'gemini-2.0-flash',
            generation_config={
                'temperature': 0.7,
                'max_output_tokens': 1500,  # Reduced for faster response
            }
        )
        
        # Concise prompt for faster generation
        prompt = f"""
Student scored {percentage:.1f}% on quiz. Weak in: {', '.join(weak_areas)}

Content (from PDF):
{quiz_content[:2000]}

Generate 2-3 CONCISE study recommendations. JSON only:
{{
  "recommendations": [
    {{
      "topic": "Topic from PDF",
      "description": "What to review (1 sentence)",
      "key_points": ["Point 1", "Point 2", "Point 3"],
      "study_tips": "How to study this (1 sentence)"
    }}
  ]
}}
"""
        
        response = model.generate_content(prompt)
        response_text = response.text.strip().replace('```json', '').replace('```', '')
        result = json.loads(response_text)
        
        return result.get('recommendations', [])[:3]  # Limit to 3 for speed
        
    except Exception as e:
        print(f" Error generating recommendations: {e}")
        return []


def student_quiz_results(request, quiz_id):
    """Enhanced results page with AI recommendations"""
    
    quiz = get_object_or_404(StudentQuiz, id=quiz_id)
    questions = QuizQuestion.objects.filter(quiz=quiz).order_by('question_number')
    
    # Calculate statistics
    total_questions = questions.count()
    score = sum(1 for q in questions if q.is_correct)
    percentage = (score / total_questions * 100) if total_questions > 0 else 0
    
    # Time formatting
    if quiz.completed_at and quiz.started_at:
        time_diff = quiz.completed_at - quiz.started_at
        minutes = int(time_diff.total_seconds() / 60)
        seconds = int(time_diff.total_seconds() % 60)
        time_taken_formatted = f"{minutes}:{seconds:02d}"
    else:
        time_taken_formatted = "N/A"
    
    # Analyze performance by question type
    type_stats = {}
    for q in questions:
        qtype = q.question_type
        if qtype not in type_stats:
            type_stats[qtype] = {'correct': 0, 'total': 0}
        type_stats[qtype]['total'] += 1
        if q.is_correct:
            type_stats[qtype]['correct'] += 1
    
    # Identify weak and strong areas
    weak_areas = []
    strong_areas = []
    for qtype, stats in type_stats.items():
        accuracy = (stats['correct'] / stats['total'] * 100) if stats['total'] > 0 else 0
        type_name = qtype.replace('_', ' ').title()
        if accuracy < 60:
            weak_areas.append(type_name)
        elif accuracy >= 80:
            strong_areas.append(type_name)
    
    # Content mastery (how well they understood the material)
    mastery_level = percentage
    
    # Generate AI recommendations based on weak areas and PDF content
    recommendations = []
    if weak_areas and percentage < 90:
        print(f" Generating AI recommendations for weak areas: {weak_areas}")
        ai_recommendations = get_study_recommendations(weak_areas, percentage, quiz.extracted_text)
        
        for rec in ai_recommendations:
            recommendations.append({
                'title': rec.get('topic', 'Study Recommendation'),
                'description': rec.get('description', 'Focus on improving this area'),
                'key_points': rec.get('key_points', []),
                'study_tips': rec.get('study_tips', '')
            })
    
    # Prepare questions review
    questions_review = []
    for q in questions:
        questions_review.append({
            'number': q.question_number,
            'type': q.question_type,
            'question_text': q.question_text,
            'student_answer': q.student_answer or 'Not answered',
            'correct_answer': q.correct_answer,
            'status': 'correct' if q.is_correct else 'incorrect',
            'reference_text': q.reference_text,
            'page_number': q.page_number
        })
    
    context = {
        'quiz': quiz,
        'score': score,
        'total_questions': total_questions,
        'percentage': percentage,
        'time_taken_formatted': time_taken_formatted,
        'mastery_level': mastery_level,
        'weak_areas': weak_areas,
        'strong_areas': strong_areas,
        'recommendations': recommendations,
        'questions_review': questions_review,
    }
    
    return render(request, 'predictor/student/quiz_results.html', context)

def landing(request):
    return render(request, 'predictor/landing.html')

def student_portal(request):
    """Student entry point - now redirects to entry page if not in session"""
    # Check if student is already logged in (via session)
    if 'student_usn' in request.session and 'student_college' in request.session:
        return render(request, 'predictor/student/home.html', {
            'student_usn': request.session['student_usn'],
            'student_college': request.session['student_college']
        })
    else:
        # Redirect to student entry page
        return redirect('student_entry')

def search_colleges_api(request):
    """Search for colleges using Gemini AI - returns JSON"""
    query = request.GET.get('q', '').strip()
    
    if not query or len(query) < 2:
        return JsonResponse({'colleges': []})
    
    # Base list of colleges (fallback)
    base_colleges = [
        # Karnataka - Bangalore
        "RV College of Engineering, Bangalore",
        "BMS College of Engineering, Bangalore", 
        "PES University, Bangalore",
        "MSRIT - Ramaiah Institute of Technology, Bangalore",
        "NMIT - Nitte Meenakshi Institute of Technology, Bangalore",
        "Dayananda Sagar College of Engineering, Bangalore",
        "CMR Institute of Technology, Bangalore",
        "BNM Institute of Technology, Bangalore",
        "RNS Institute of Technology, Bangalore",
        "Presidency University, Bangalore",
        
        # Karnataka - Other Cities
        "NIE Institute of Technology, Mysore",
        "SJCE - JSS Science and Technology University, Mysuru",
        "Manipal Institute of Technology, Manipal",
        "Shri Madhwa Vadiraja Institute of Technology, Bantakal",
        "MIT - Manipal Institute of Technology, Moodalkatte",
        "Alvas Institute of Engineering and Technology, Moodabidri",
        "Canara Engineering College, Mangalore",
        "NMAM Institute of Technology, Nitte",
        "Sahyadri College of Engineering, Mangalore",
        
        # Delhi
        "IIT Delhi - Indian Institute of Technology Delhi",
        "DTU - Delhi Technological University",
        "NSIT - Netaji Subhas Institute of Technology, Delhi",
        
        # Mumbai/Maharashtra  
        "IIT Bombay - Indian Institute of Technology Bombay",
        "VJTI - Veermata Jijabai Technological Institute, Mumbai",
        "COEP - College of Engineering Pune",
        
        # Tamil Nadu
        "IIT Madras - Indian Institute of Technology Madras",
        "Anna University, Chennai",
        "PSG College of Technology, Coimbatore",
        
        # Other Major Colleges
        "BITS Pilani - Birla Institute of Technology and Science Pilani",
        "VIT Vellore - Vellore Institute of Technology",
        "NIT Trichy - National Institute of Technology Tiruchirappalli",
        "NIT Surathkal - National Institute of Technology Karnataka"
    ]
    
    # Filter base colleges
    matches = [c for c in base_colleges if query.lower() in c.lower()]
    
    # If Gemini API is available, enhance with AI search
    if genai and GEMINI_API_KEY and len(matches) < 5:
        try:
            model = genai.GenerativeModel('gemini-2.0-flash')
            prompt = f"""List 10 engineering colleges in India that match the search query: "{query}"
            
Include colleges from Karnataka (especially coastal Karnataka), other states, and variations of the searched name.
Return ONLY college names with city, one per line, no numbering or extra text.
Format: College Name, City"""
            
            response = model.generate_content(prompt)
            ai_colleges = [line.strip() for line in response.text.strip().split('\n') if line.strip()]
            
            # Add AI results to matches (avoid duplicates)
            for college in ai_colleges[:10]:
                if college and college not in matches:
                    matches.append(college)
        except Exception as e:
            print(f"Gemini API error: {e}")
    
    # Return top 15 matches
    return JsonResponse({'colleges': matches[:15]})

def student_entry(request):
    """New student entry page with USN and college - STRICT ACCESS CONTROL"""
    from .models import StudentRecord, StudentNotification
    
    if request.method == 'POST':
        usn = request.POST.get('usn', '').strip().upper()
        college = request.POST.get('college', '').strip()
        
        if not usn or not college:
            return render(request, 'predictor/student/student_entry.html', {
                'error': 'Please fill in all required fields',
                'usn': usn,
                'college': college
            })
        
        # STRICT CHECK: Only allow students who exist in college portal database
        try:
            student = StudentRecord.objects.get(student_id=usn)
            # Student exists in college portal - GRANT ACCESS
            request.session['student_usn'] = usn
            request.session['student_college'] = college
            request.session['student_name'] = student.name
            print(f" Access granted to student: {usn} - {student.name}")
            return redirect('student_portal')
            
        except StudentRecord.DoesNotExist:
            # Student NOT in college portal database - DENY ACCESS
            # Create notification for tracking
            StudentNotification.objects.get_or_create(
                usn=usn,
                college_name=college,
                defaults={'is_resolved': False}
            )
            print(f" Access denied - USN not found in college portal: {usn}")
            
            # Show error message asking to contact college
            return render(request, 'predictor/student/student_entry.html', {
                'error': 'Your USN is not registered in the college portal. Please ask your college administrator to add your details to access the dashboard.',
                'usn': usn,
                'college': college,
                'access_denied': True
            })
    
    return render(request, 'predictor/student/student_entry.html')

def student_predict(request):
    from .models import StudentRecord, StudentPrediction
    
    if request.method == 'POST':
        # Get student_id from session first, fallback to POST data
        student_id = request.session.get('student_usn') or request.POST.get('student_id')
        
        if not student_id:
            return render(request, 'predictor/student/predict.html', {'error': 'Session expired. Please log in again.'})
        
        try:
            student = StudentRecord.objects.get(student_id=student_id)
            communication = float(request.POST.get('communication_skills', 5))
            technical = float(request.POST.get('technical_skills', 5))
            aptitude = float(request.POST.get('aptitude_score', 50))
            projects = int(request.POST.get('projects', 0))
            internships = int(request.POST.get('internships', 0))
            certifications = int(request.POST.get('certifications', 0))
            
            # Calculate placement probability
            score = student.cgpa * 10 + communication * 5 + technical * 5 + aptitude * 0.3 + projects * 2 + internships * 3 + certifications * 2
            probability = min(score / 2, 100)
            
            prediction = StudentPrediction.objects.create(
                student=student, cgpa=student.cgpa, backlogs=student.total_backlogs,
                communication_skills=communication, technical_skills=technical, aptitude_score=aptitude,
                projects_completed=projects, internships=internships, certifications=certifications,
                placement_probability=probability, prediction='Placed' if probability > 60 else 'Not Placed',
                confidence_score=85.0, recommendations='Improve skills.'
            )
            return render(request, 'predictor/student/prediction_result.html', {'student': student, 'prediction': prediction})
        except StudentRecord.DoesNotExist:
            return render(request, 'predictor/student/predict.html', {'error': 'Student ID not found. Your college needs to add your record.'})
    return render(request, 'predictor/student/predict.html')

def student_results(request):
    from .models import StudentQuiz, StudentPrediction
    
    # Get the logged-in student's USN from session
    student_usn = request.session.get('student_usn')
    student_name = request.session.get('student_name', 'Student')
    
    if not student_usn:
        # Redirect to login if no session found
        return redirect('student_entry')
    
    # Get quiz results - filter by student_id and completed status
    quizzes = StudentQuiz.objects.filter(
        student_id=student_usn,
        status='completed'
    ).order_by('-completed_at')
    
    # Get placement predictions - filter by student's student_id
    predictions = StudentPrediction.objects.filter(
        student__student_id=student_usn
    ).order_by('-predicted_at')
    
    return render(request, 'predictor/student/results.html', {
        'quizzes': quizzes,
        'predictions': predictions,
        'student_usn': student_usn,
        'student_name': student_name
    })


# ==================== STUDENT TRAINING SESSIONS ====================

def student_sessions(request):
    """View all training sessions - student can enroll/unenroll"""
    # Get the logged-in student's USN from session
    student_usn = request.session.get('student_usn')
    
    if not student_usn:
        # Redirect to login if no session found
        messages.warning(request, 'Please enter your Student ID to access sessions.')
        return redirect('student_entry')
    
    # Get student record
    try:
        student = StudentRecord.objects.get(student_id=student_usn)
    except StudentRecord.DoesNotExist:
        messages.error(request, 'Student record not found.')
        return redirect('student_entry')
    
    # Get all active training sessions
    all_sessions = TrainingSession.objects.filter(
        is_active=True,
        scheduled_date__gte=timezone.now().date()
    ).order_by('scheduled_date', 'scheduled_time')
    
    # Get sessions student is enrolled in
    enrolled_session_ids = student.training_sessions.values_list('id', flat=True)
    
    # Get recommendations for this student
    recommendations = SessionRecommendation.objects.filter(
        student=student
    ).select_related('session')
    
    recommended_session_ids = [rec.session_id for rec in recommendations]
    
    # Build session data with enrollment status
    sessions_data = []
    for session in all_sessions:
        enrolled_count = session.registered_students.count()
        is_full = enrolled_count >= session.max_students
        is_enrolled = session.id in enrolled_session_ids
        is_recommended = session.id in recommended_session_ids
        
        # Get recommendation details if exists
        recommendation = None
        if is_recommended:
            recommendation = recommendations.filter(session=session).first()
        
        sessions_data.append({
            'session': session,
            'enrolled_count': enrolled_count,
            'seats_left': session.max_students - enrolled_count,
            'is_full': is_full,
            'is_enrolled': is_enrolled,
            'is_recommended': is_recommended,
            'recommendation': recommendation,
            'resource_list': session.resource_links.split(',') if session.resource_links else []
        })
    
    context = {
        'sessions': sessions_data,
        'student_usn': student_usn,
        'student_name': request.session.get('student_name', 'Student'),
    }
    
    return render(request, 'predictor/student/sessions.html', context)


def enroll_session(request, session_id):
    """Enroll student in a training session"""
    if request.method != 'POST':
        return JsonResponse({'success': False, 'message': 'Invalid request method'})
    
    student_usn = request.session.get('student_usn')
    if not student_usn:
        return JsonResponse({'success': False, 'message': 'Please login first'})
    
    try:
        student = StudentRecord.objects.get(student_id=student_usn)
        session = get_object_or_404(TrainingSession, id=session_id, is_active=True)
        
        # Check if already enrolled
        if student in session.registered_students.all():
            return JsonResponse({'success': False, 'message': 'Already enrolled in this session'})
        
        # Check if session is full
        if session.registered_students.count() >= session.max_students:
            return JsonResponse({'success': False, 'message': 'Session is full'})
        
        # Enroll the student
        session.registered_students.add(student)
        
        # Mark recommendation as registered if exists
        SessionRecommendation.objects.filter(
            student=student,
            session=session
        ).update(is_registered=True)
        
        enrolled_count = session.registered_students.count()
        seats_left = session.max_students - enrolled_count
        
        return JsonResponse({
            'success': True,
            'message': f'Successfully enrolled in {session.title}',
            'enrolled_count': enrolled_count,
            'seats_left': seats_left
        })
        
    except StudentRecord.DoesNotExist:
        return JsonResponse({'success': False, 'message': 'Student record not found'})
    except Exception as e:
        return JsonResponse({'success': False, 'message': str(e)})


def unenroll_session(request, session_id):
    """Unenroll student from a training session"""
    if request.method != 'POST':
        return JsonResponse({'success': False, 'message': 'Invalid request method'})
    
    student_usn = request.session.get('student_usn')
    if not student_usn:
        return JsonResponse({'success': False, 'message': 'Please login first'})
    
    try:
        student = StudentRecord.objects.get(student_id=student_usn)
        session = get_object_or_404(TrainingSession, id=session_id, is_active=True)
        
        # Check if enrolled
        if student not in session.registered_students.all():
            return JsonResponse({'success': False, 'message': 'Not enrolled in this session'})
        
        # Check if session is too soon (e.g., within 24 hours)
        time_until_session = session.scheduled_date - timezone.now().date()
        if time_until_session.days < 1:
            return JsonResponse({
                'success': False,
                'message': 'Cannot unenroll less than 24 hours before session'
            })
        
        # Unenroll the student
        session.registered_students.remove(student)
        
        # Mark recommendation as not registered if exists
        SessionRecommendation.objects.filter(
            student=student,
            session=session
        ).update(is_registered=False)
        
        enrolled_count = session.registered_students.count()
        seats_left = session.max_students - enrolled_count
        
        return JsonResponse({
            'success': True,
            'message': f'Successfully unenrolled from {session.title}',
            'enrolled_count': enrolled_count,
            'seats_left': seats_left
        })
        
    except StudentRecord.DoesNotExist:
        return JsonResponse({'success': False, 'message': 'Student record not found'})
    except Exception as e:
        return JsonResponse({'success': False, 'message': str(e)})


